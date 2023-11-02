import os
from flask import Blueprint, render_template, request, send_file, flash
from io import BytesIO
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image as PILImage

ppttopdf_bp = Blueprint('ppttopdf', __name__)

@ppttopdf_bp.route('/ppttopdf', methods=['GET', 'POST'])
def ppttopdfindex():
    if request.method == 'POST':
        if 'pptx_file' not in request.files:
            flash('No file part', 'error')
        else:
            pptx_file = request.files['pptx_file']
            if pptx_file.filename == '':
                flash('No selected file', 'error')
            else:
                pptx_bytes = pptx_file.read()
                pdf_bytes, error_message = convert_pptx_to_pdf(pptx_bytes)
                if pdf_bytes:
                    return send_file(
                        BytesIO(pdf_bytes),
                        as_attachment=True,
                        download_name='converted.pdf',
                        mimetype='application/pdf'
                    )
                else:
                    flash(error_message, 'error')

    return render_template('ppttopdf.html')

def convert_pptx_to_pdf(pptx_bytes):
    try:
        prs = Presentation(BytesIO(pptx_bytes))
        pdf_buffer = BytesIO()
        pdf_canvas = canvas.Canvas(pdf_buffer, pagesize=letter)

        for i, slide in enumerate(prs.slides):
            image = None
            for shape in slide.shapes:
                if shape.shape_type == 13:  # 13 represents an image
                    if hasattr(shape, 'image') and shape.image is not None:
                        image_bytes = shape.image.blob
                        image = PILImage.open(BytesIO(image_bytes))
                        break

            if image:
                image.thumbnail(letter)
                slide_image_path = f"slide_{i}.png"
                image.save(slide_image_path)

                if i > 0:
                    pdf_canvas.showPage()

                pdf_canvas.drawImage(slide_image_path, 0, 0, width=letter[0], height=letter[1])
                os.remove(slide_image_path)

        pdf_canvas.save()
        pdf_bytes = pdf_buffer.getvalue()
        return pdf_bytes, None
    except Exception as e:
        error_message = f"Conversion error: {str(e)}"
        print(error_message)
        return None, error_message
