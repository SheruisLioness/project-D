import os
import zipfile
from flask import Blueprint, render_template, request, send_file, flash
from pptx import Presentation
from io import BytesIO
from PIL import Image

ppttoimages_bp = Blueprint('ppttoimages', __name__)

def pptx_to_images(pptx_bytes):
    images = []
    prs = Presentation(BytesIO(pptx_bytes))
    
    for i, slide in enumerate(prs.slides):
        image = slide_to_image(slide)
        images.append(image)

    return images

def slide_to_image(slide):
    slide_width = 1024  # Adjust this width as needed
    slide_height = 768  # Adjust this height as needed

    slide_image = Image.new("RGB", (slide_width, slide_height), (255, 255, 255))

    for shape in slide.shapes:
        if shape.shape_type == 13:  # Image shape
            if shape.image:
                with BytesIO(shape.image.blob) as image_stream:
                    img = Image.open(image_stream)
                    img = img.convert("RGB")
                    img.thumbnail((slide_width, slide_height))
                    slide_image.paste(img, (0, 0))

    return slide_image

@ppttoimages_bp.route('/ppttoimg',  methods=['GET', 'POST'])
def ppttoimgindex():
    if request.method == 'POST':
        if 'pptx_file' not in request.files:
            flash('No file part', 'error')
            return render_template('ppttoimages.html')

        pptx_file = request.files['pptx_file']
        if pptx_file.filename == '':
            flash('No selected file', 'error')
            return render_template('ppttoimages.html')

        pptx_bytes = pptx_file.read()
        images = pptx_to_images(pptx_bytes)
        
        zip_bytes = create_zip(images)
        
        if zip_bytes:
            return send_file(
                BytesIO(zip_bytes),
                as_attachment=True,
                download_name='pptx_images.zip',
                mimetype='application/zip'
            )
        else:
            flash('Conversion failed', 'error')
            return render_template('ppttoimages.html')
    else:
        return render_template('ppttoimages.html')

def create_zip(images):
    try:
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for i, image in enumerate(images):
                image_byte_io = BytesIO()
                image.save(image_byte_io, format='PNG')
                zipf.writestr(f'image_{i}.png', image_byte_io.getvalue())

        return zip_buffer.getvalue()
    except Exception as e:
        print(f"ZIP creation error: {str(e)}")
        return None
