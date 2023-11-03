import pythoncom
from flask import Flask, Blueprint, redirect, render_template, request, send_file, url_for
import os
from docx2pdf import convert as docx_to_pdf
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import datetime

from email_utils import send_email

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'

docxtopptx_blueprint = Blueprint('docxtopptx', __name__)

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def convert_to_pdf(docx_path):
    pdf_path = os.path.splitext(docx_path)[0] + '.pdf'
    docx_to_pdf(docx_path, pdf_path)
    return pdf_path

def extract_images(pdf_path):
    images = convert_from_path(pdf_path)
    image_paths = []

    for i, image in enumerate(images):
        image_path = os.path.join(UPLOAD_FOLDER, f'page_{i}.png')
        image.save(image_path, 'PNG')
        image_paths.append(image_path)

    return image_paths

def extract_text(docx_path):
    from docx import Document
    doc = Document(docx_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return text

def create_pptx(text, images, output_path):
    pythoncom.CoInitialize()  # Initialize the COM library
    prs = Presentation()

    for i, slide_content in enumerate(text):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, content = slide.placeholders
        title.text = ''
        content.text = slide_content

        if i < len(images) and images[i] is not None:
            img_path = images[i]
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5.5)
            slide.shapes.add_picture(img_path, left, top, width, height)

    prs.save(output_path)

@docxtopptx_blueprint.route('/docxtopptx', methods=['GET', 'POST'])
def docxtopptx_index():
    pythoncom.CoInitialize()  # Initialize the COM library
    if request.method == 'POST':
        if 'docx_file' not in request.files:
            return render_template('docxtopptx.html', error='No file part')

        docx_file = request.files['docx_file']

        if docx_file.filename == '':
            return render_template('docxtopptx.html', error='No selected file')

        if docx_file:
            docx_path = os.path.join(UPLOAD_FOLDER, docx_file.filename)
            docx_file.save(docx_path)

            pdf_path = convert_to_pdf(docx_path)
            images = extract_images(pdf_path)
            text = extract_text(docx_path)

            now = datetime.datetime.now()
            pptx_filename = now.strftime('%Y-%m-%d_%H-%M-%S') + '.pptx'
            pptx_path = os.path.join(UPLOAD_FOLDER, pptx_filename)

            create_pptx(text, images, pptx_path)

            for image_path in images:
                os.remove(image_path)
            os.remove(pdf_path)

            return redirect(url_for('docxtopptx.download_pptx', pptx_filename=pptx_filename, _external=True, _scheme='http'))

    return render_template('docxtopptx.html', error=None)

@docxtopptx_blueprint.route('/download_pptx/<pptx_filename>', methods=['GET','POST'])
def download_pptx(pptx_filename):
    action = request.form.get('action')
    pptx_path = os.path.join(UPLOAD_FOLDER, pptx_filename)

    if action == 'send':
        recipient_email = request.form.get('email')
        smtp_username = 'your_email@gmail.com'  # Replace with your Gmail email address
        smtp_password = 'your_app_password'  # Replace with your app password
        send_email(pptx_path, recipient_email, smtp_username, smtp_password)
        return redirect(url_for('docxtopptx.docxtopptx_index'))
    elif action == 'download':
        return send_file(pptx_path, as_attachment=True, attachment_filename=pptx_filename)
    else:
        return "Invalid action"

app.register_blueprint(docxtopptx_blueprint)

if __name__ == '__main__':
    app.run(debug=True)
