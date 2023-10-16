import os
from flask import Blueprint, render_template, request, send_file, flash
from pptx import Presentation
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image

ppttopdf_bp = Blueprint('ppttopdf', __name__)

@ppttopdf_bp.route('/ppttopdf', methods=['GET', 'POST'])
def ppttopdfindex():
    if request.method == 'POST':
        if 'pptx_file' not in request.files:
            flash('No file part', 'error')
        else:
            pptx_file = request.files['pptx_file']
            if pptx_file.filename == '':
                flash('No selected file', 'error')
            else:
                pptx_bytes = pptx_file.read()
                pdf_bytes, error_message = convert_pptx_to_pdf(pptx_bytes)
                if pdf_bytes:
                    return send_file(
                        BytesIO(pdf_bytes),
                        as_attachment=True,
                        download_name='converted.pdf',
                        mimetype='application/pdf'
                    )
                else:
                    flash(error_message, 'error')

    return render_template('ppttopdf.html')

def convert_pptx_to_pdf(pptx_bytes):
    try:
        prs = Presentation(BytesIO(pptx_bytes))
        pdf_buffer = BytesIO()

        # Create a PDF canvas
        pdf_canvas = canvas.Canvas(pdf_buffer, pagesize=letter)

        for slide in prs.slides:
            slide_width, slide_height = letter
            image_stream = BytesIO()

            slide_image = Image.new("RGB", (int(slide_width), int(slide_height)), (255, 255, 255))

            for shape in slide.shapes:
                if shape.shape_type == 13:  # 13 represents an image
                    image = shape.image
                    image_bytes = image.blob
                    image_stream.write(image_bytes)
                    image_stream.seek(0)

                    img = Image.open(image_stream)
                    img.thumbnail((int(slide_width), int(slide_height)))  # corrected line
                    pdf_canvas.drawInlineImage(img, 0, 0, width=letter[0], height=letter[1])

            slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            pdf_canvas.drawString(100, 100, slide_text)  # Adjust the position as needed
            pdf_canvas.showPage()

        pdf_canvas.save()
        pdf_bytes = pdf_buffer.getvalue()
        return pdf_bytes, None
    except Exception as e:
        error_message = f"Conversion error: {str(e)}"
        print(error_message)
        return None, error_message
