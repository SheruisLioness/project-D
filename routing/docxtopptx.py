from flask import Flask, Blueprint, render_template, request, send_file
import os
from docx2pdf import convert as docx_to_pdf
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import datetime

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'

docxtopptx_blueprint = Blueprint('docxtopptx', __name__)

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def convert_to_pdf(docx_path):
    pdf_path = os.path.splitext(docx_path)[0] + '.pdf'
    docx_to_pdf(docx_path, pdf_path)
    return pdf_path

def extract_images(pdf_path):
    images = convert_from_path(pdf_path)
    image_paths = []

    for i, image in enumerate(images):
        image_path = os.path.join(UPLOAD_FOLDER, f'page_{i}.png')
        image.save(image_path, 'PNG')
        image_paths.append(image_path)

    return image_paths

def create_pptx(images, output_path):
    prs = Presentation()

    for image_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        left = top = Inches(1)
        width = height = Inches(8)

        slide.shapes.add_picture(image_path, left, top, width, height)

    prs.save(output_path)

@docxtopptx_blueprint.route('/docxtopptx', methods=['GET', 'POST'])
def docxtopptx_index():
    if request.method == 'POST':
        if 'docx_file' not in request.files:
            return render_template('docxtopptx.html', error='No file part')

        docx_file = request.files['docx_file']
        
        if docx_file.filename == '':
            return render_template('docxtopptx.html', error='No selected file')

        if docx_file:
            docx_path = os.path.join(UPLOAD_FOLDER, docx_file.filename)
            docx_file.save(docx_path)

            pdf_path = convert_to_pdf(docx_path)
            images = extract_images(pdf_path)
            
            now = datetime.datetime.now()
            pptx_filename = now.strftime('%Y-%m-%d_%H-%M-%S') + '.pptx'
            pptx_path = os.path.join(UPLOAD_FOLDER, pptx_filename)
            
            create_pptx(images, pptx_path)
            
            for image_path in images:       
                os.remove(image_path)
            os.remove(pdf_path)
            
            return render_template('download_pptx.html', pptx_filename=pptx_filename)
    
    return render_template('docxtopptx.html', error=None)

@docxtopptx_blueprint.route('/download_pptx/<pptx_filename>')
def download_pptx(pptx_filename):
    pptx_path = os.path.join(UPLOAD_FOLDER, pptx_filename)
    return send_file(pptx_path, as_attachment=True)

app.register_blueprint(docxtopptx_blueprint)

if __name__ == '__main__':
    app.run(debug=True)
