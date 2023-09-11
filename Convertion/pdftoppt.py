import os
from flask import Flask, render_template, request, send_file
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import tempfile
import datetime  # Import the datetime module

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

def pdf_to_images(pdf_bytes):
    images = convert_from_bytes(pdf_bytes)
    return images

def create_pptx(images):
    prs = Presentation()
    
    # Get the slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    temp_dir = tempfile.mkdtemp()  # Create a temporary directory

    for i, image in enumerate(images):
        image_path = os.path.join(temp_dir, f'image_{i}.png')
        image.save(image_path, 'PNG')  # Save the image as a temporary file
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        left = 0  # Set the left position to 0 inches
        top = 0   # Set the top position to 0 inches
        pic = slide.shapes.add_picture(image_path, left, top, width=slide_width, height=slide_height)

    # Generate the filename using the current date and time
    current_datetime = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
    pptx_filename = f'output_{current_datetime}.pptx'
    pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], pptx_filename)
    prs.save(pptx_path)

    # Clean up temporary directory
    for filename in os.listdir(temp_dir):
        file_path = os.path.join(temp_dir, filename)
        os.remove(file_path)
    os.rmdir(temp_dir)

    return pptx_path

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        if 'pdf_file' not in request.files:
            return render_template('pdftoppt.html', error='No file part')

        pdf_file = request.files['pdf_file']
        
        if pdf_file.filename == '':
            return render_template('pdftoppt.html', error='No selected file')

        if pdf_file:
            pdf_bytes = pdf_file.read()
            images = pdf_to_images(pdf_bytes)
            pptx_path = create_pptx(images)
            return send_file(pptx_path, as_attachment=True, download_name=os.path.basename(pptx_path))
    
    return render_template('pdftoppt.html', error=None)

if __name__ == '__main__':
    app.run(debug=True)
