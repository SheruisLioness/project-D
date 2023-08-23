import os
from flask import Flask, render_template, request, send_file
from docx import Document
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

def convert_to_pptx(input_path, output_path):
    doc = Document(input_path)
    prs = Presentation()

    slide = None

    for para in doc.paragraphs:
        if slide is None:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Choose a slide layout (e.g., Title Slide)

        left = top = Inches(1)
        width = height = Inches(8)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        frame = textbox.text_frame

        for run in para.runs:
            frame.text += run.text + ' '

        if para == doc.paragraphs[-1] or para.style.name == 'Heading 1':  # Start a new slide for each heading or at the end
            slide = None

    prs.save(output_path)

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        if 'docx_file' not in request.files:
            return render_template('doctopptx.html', error='No file part')

        docx_file = request.files['docx_file']
        
        if docx_file.filename == '':
            return render_template('doctopptx.html', error='No selected file')

        if docx_file:
            docx_path = os.path.join(app.config['UPLOAD_FOLDER'], docx_file.filename)
            docx_file.save(docx_path)

            pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], 'output.pptx')
            convert_to_pptx(docx_path, pptx_path)
            
            return render_template('download_pptx.html', pptx_path='output.pptx')
    
    return render_template('doctopptx.html', error=None)

@app.route('/download_pptx')
def download_pptx():
    pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], 'output.pptx')
    return send_file(pptx_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
